#!/usr/bin/env python3
"""Builds the board deck for the infrastructure cost programme.

Mirrors Infrastructure/infrastructure-cost-plans-summary.html (v4.0) — same
figures, same section order, recommendation last. Re-run after editing the
report so the two stay in step:  python3 build_cost_deck.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x12, 0x31, 0x4F)
NAVY2 = RGBColor(0x1C, 0x4A, 0x78)
INK = RGBColor(0x0F, 0x22, 0x33)
BODY = RGBColor(0x33, 0x42, 0x4F)
MUTED = RGBColor(0x5F, 0x71, 0x83)
RULE = RGBColor(0xD8, 0xDE, 0xE5)
TINT = RGBColor(0xF4, 0xF7, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
POS = RGBColor(0x15, 0x65, 0x3F)
NEG = RGBColor(0x9A, 0x4B, 0x18)
ALERT = RGBColor(0xA4, 0x26, 0x2C)

HEAD_FONT = "Segoe UI Semibold"
BODY_FONT = "Segoe UI"
MONO_FONT = "Consolas"

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.62)
CONTENT_W = W - 2 * MARGIN

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]

_slide_no = 0


def box(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0.75)):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = line_w
    return sh


def text(slide, x, y, w, h, runs, size=14, color=BODY, font=BODY_FONT,
         bold=False, align=PP_ALIGN.LEFT, space=6, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.15):
    """runs: str, or list of paragraphs; a paragraph is a str or list of
    (text, {opts}) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paras = runs if isinstance(runs, list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        p.line_spacing = line_spacing
        pieces = para if isinstance(para, list) else [(para, {})]
        for t, o in pieces:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(o.get("size", size))
            r.font.bold = o.get("bold", bold)
            r.font.name = o.get("font", font)
            r.font.color.rgb = o.get("color", color)
            r.font.italic = o.get("italic", False)
    return tb


def slide_frame(title, kicker=None, lede=None):
    """Standard content slide: masthead rule, section title, optional lede."""
    global _slide_no
    s = prs.slides.add_slide(BLANK)
    _slide_no += 1
    box(s, Emu(0), Emu(0), W, H, fill=WHITE)
    # header rule
    box(s, MARGIN, Inches(0.95), CONTENT_W, Pt(1.6), fill=INK)
    text(s, MARGIN, Inches(0.34), CONTENT_W, Inches(0.3),
         [[("PLUGINLIVE", {"color": NAVY, "bold": True, "size": 10.5}),
           ("     Infrastructure cost optimisation — summary of all plans", {"color": MUTED, "size": 10.5}),
           ]], size=10.5)
    if kicker:
        text(s, MARGIN, Inches(1.16), CONTENT_W, Inches(0.26), kicker,
             size=10.5, color=NAVY2, font=MONO_FONT)
    y = Inches(1.44) if kicker else Inches(1.24)
    text(s, MARGIN, y, CONTENT_W, Inches(0.5), title, size=27, color=INK,
         font=HEAD_FONT, bold=True)
    y = y + Inches(0.62)
    if lede:
        text(s, MARGIN, y, Inches(11.4), Inches(0.9), lede, size=13.5, color=BODY)
        y = y + Inches(0.30) * (1 + len(lede) // 96)
    # footer
    text(s, MARGIN, H - Inches(0.52), Inches(9), Inches(0.25),
         "Internal — Confidential · Version 4.0 · 25 August 2026", size=9, color=MUTED)
    text(s, W - MARGIN - Inches(1.0), H - Inches(0.52), Inches(1.0), Inches(0.25),
         str(_slide_no), size=9, color=MUTED, align=PP_ALIGN.RIGHT)
    return s, y + Inches(0.18)


def table(slide, x, y, w, headers, rows, col_w=None, font=11,
          head_font=9.5, row_h=Inches(0.34), head_h=Inches(0.33)):
    """Navy header, hairline rows. rows: list of list of (text, opts) or str."""
    n_rows, n_cols = len(rows) + 1, len(headers)
    total_h = head_h + row_h * len(rows)
    shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, total_h)
    tbl = shape.table
    tbl.first_row = True
    tbl.horz_banding = False
    if col_w:
        units = sum(col_w)
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Emu(int(w * cw / units))
    tbl.rows[0].height = head_h
    for r in range(1, n_rows):
        tbl.rows[r].height = row_h

    for c, htxt in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.margin_left = cell.margin_right = Inches(0.09)
        cell.margin_top = cell.margin_bottom = Inches(0.04)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT if htxt.startswith(">") else PP_ALIGN.LEFT
        r = p.add_run()
        r.text = htxt.lstrip(">").upper()
        r.font.size = Pt(head_font)
        r.font.bold = True
        r.font.name = BODY_FONT
        r.font.color.rgb = WHITE

    for ri, row in enumerate(rows, start=1):
        for ci, cellspec in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 else TINT
            cell.margin_left = cell.margin_right = Inches(0.09)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            txt, o = (cellspec, {}) if isinstance(cellspec, str) else cellspec
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if o.get("num") else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(o.get("size", font))
            r.font.bold = o.get("bold", ci == 0)
            r.font.name = MONO_FONT if o.get("num") else BODY_FONT
            r.font.color.rgb = o.get("color", INK if ci == 0 or o.get("num") else BODY)
    return shape


def bullets(slide, x, y, w, items, size=13, h=Inches(3.0), dash_color=NAVY2,
            space=Pt(7)):
    """One textbox, one paragraph per item, hanging indent — PowerPoint flows
    the wrapping so blocks can never overlap."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = space
        p.line_spacing = 1.12
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", "165100")     # 0.18" indent
        pPr.set("indent", "-165100")  # hanging
        r = p.add_run()
        r.text = "—  "
        r.font.size = Pt(size)
        r.font.name = BODY_FONT
        r.font.color.rgb = dash_color
        pieces = item if isinstance(item, list) else [(item, {})]
        for t, o in pieces:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(o.get("size", size))
            r.font.bold = o.get("bold", False)
            r.font.name = o.get("font", BODY_FONT)
            r.font.color.rgb = o.get("color", BODY)
    return tb


def _plain(para):
    if isinstance(para, str):
        return para
    return "".join(t for t, _ in para)


def kpi_row(slide, y, tiles, h=Inches(1.18)):
    n = len(tiles)
    gap = Inches(0.14)
    w = int((CONTENT_W - gap * (n - 1)) / n)
    for i, (label, value, note, hi) in enumerate(tiles):
        x = MARGIN + i * (w + gap)
        box(slide, x, y, Emu(w), h, fill=TINT if hi else WHITE, line=RULE)
        text(slide, x + Inches(0.16), y + Inches(0.14), Emu(w) - Inches(0.3), Inches(0.22),
             label.upper(), size=9, color=MUTED, font=MONO_FONT)
        text(slide, x + Inches(0.16), y + Inches(0.4), Emu(w) - Inches(0.3), Inches(0.44),
             value, size=25, color=NAVY, font=HEAD_FONT, bold=True)
        text(slide, x + Inches(0.16), y + Inches(0.86), Emu(w) - Inches(0.3), Inches(0.22),
             note, size=9.5, color=MUTED)


# ───────────────────────────── 1 · title ─────────────────────────────
s = prs.slides.add_slide(BLANK)
box(s, Emu(0), Emu(0), W, H, fill=WHITE)
box(s, Emu(0), Emu(0), Inches(0.34), H, fill=NAVY)
text(s, Inches(1.15), Inches(1.15), Inches(3), Inches(0.3), "PLUGINLIVE",
     size=13, color=NAVY, bold=True, font=HEAD_FONT)
box(s, Inches(1.15), Inches(1.62), Inches(1.4), Pt(2.2), fill=NAVY)
text(s, Inches(1.15), Inches(2.1), Inches(9.6), Inches(2.0),
     "Infrastructure cost\noptimisation", size=52, color=INK, font=HEAD_FONT,
     bold=True, line_spacing=0.98)
text(s, Inches(1.15), Inches(4.05), Inches(9.3), Inches(0.8),
     "Summary of all plans — Phase 0, Plan A, Plan B1, Plan B2, Plan C",
     size=17, color=NAVY2)
box(s, Inches(1.15), Inches(4.75), Inches(9.3), Pt(1), fill=RULE)
meta = [("Version", "4.0 · consolidated"), ("Date", "25 August 2026"),
        ("Prepared by", "Infrastructure / FinOps"), ("Decision owner", "CTO"),
        ("Status", "For approval")]
for i, (k, v) in enumerate(meta):
    x = Inches(1.15) + i * Inches(1.92)
    text(s, x, Inches(5.05), Inches(1.8), Inches(0.2), k.upper(), size=8.5,
         color=MUTED, font=MONO_FONT)
    text(s, x, Inches(5.3), Inches(1.8), Inches(0.24), v, size=11.5, color=INK, bold=True)
box(s, Inches(1.15), Inches(6.25), Inches(2.5), Inches(0.34), line=ALERT)
text(s, Inches(1.15), Inches(6.33), Inches(2.5), Inches(0.24), "INTERNAL — CONFIDENTIAL",
     size=9, color=ALERT, font=MONO_FONT, align=PP_ALIGN.CENTER)
_slide_no = 1

# ───────────────────────────── 2 · position ─────────────────────────────
s, y = slide_frame("The position today", kicker="SECTION 1 · PURPOSE AND FINANCIAL POSITION",
                   lede="Two hard constraints: everything stays in an India region, and performance must be no worse than today. All four options satisfy both — they differ by roughly ₹5,000 to ₹23,000 per month, and by an order of magnitude in migration risk.")
kpi_row(s, y + Inches(0.06), [
    ("Original run-rate", "₹1,27,600", "before Phase 0", False),
    ("Already banked", "₹44,000", "Phase 0, complete", False),
    ("Current run-rate", "₹84,800", "baseline for this paper", False),
    ("Plan A target", "₹21,000", "−75% versus today", True),
    ("Alternatives", "₹26k–44k", "B2 · C · B1", False),
])
table(s, MARGIN, y + Inches(1.48), CONTENT_W,
      ["Stage / option", ">₹ per month", ">Versus today", "Status"],
      [["Original estate", ("1,27,600", {"num": True}), ("—", {"num": True}), "Superseded"],
       ["Phase 0 — duplicate PG14 estate removed", ("−44,000", {"num": True, "color": POS}), ("—", {"num": True}), "Complete"],
       ["Current run-rate", ("84,800", {"num": True}), ("baseline", {"num": True}), "Live"],
       ["Plan A — re-shape on Oracle Cloud", ("≈ 21,000", {"num": True}), ("−75%", {"num": True, "color": POS}), "Under appraisal"],
       ["Plan B1 — Cloud Run, lift and shift", ("≈ 44,000", {"num": True}), ("−48%", {"num": True, "color": POS}), "Under appraisal"],
       ["Plan B2 — Cloud Run, after queue rewrite", ("≈ 26,000", {"num": True}), ("−69%", {"num": True, "color": POS}), "Under appraisal"],
       ["Plan C — Azure Container Apps, self-managed DB", ("≈ 41,300", {"num": True}), ("−51%", {"num": True, "color": POS}), "Under appraisal"]],
      col_w=[46, 16, 16, 22], row_h=Inches(0.30))

# ───────────────────────────── 3 · four options ─────────────────────────────
s, y = slide_frame("Four costed options", kicker="SECTION 2 · OPTIONS APPRAISAL")
cards = [
    ("Option A", "Plan A — Oracle Cloud", "₹21,000", "7 weeks · moderate effort",
     ["−75% versus today", "₹8,488 saved in week 0", "No data migration", "Reversible at every step"], []),
    ("Option B1", "Plan B1 — Cloud Run", "₹44,000", "10 weeks · high effort",
     ["−48% versus today", "Native scale-to-zero"],
     ["Workers held warm — ₹11,500/mo", "10 GB DB + 160 GB media move"]),
    ("Option B2", "Plan B2 — + queue rewrite", "₹26,000", "14 weeks · very high effort",
     ["−69% versus today", "Workers scale to zero"],
     ["Requires BullMQ rewrite", "Still ₹5,000/mo worse than A"]),
    ("Option C", "Plan C — Azure", "₹41,300", "9–12 weeks · high effort",
     ["−51%; ₹27,000 fully optimised", "KEDA built in — no rewrite"],
     ["Same migration as Plan B", "Memory floor near ₹30,000"]),
]
cw = int((CONTENT_W - Inches(0.16) * 3) / 4)
for i, (chip, name, price, per, pros, cons) in enumerate(cards):
    x = MARGIN + i * (cw + Inches(0.16))
    box(s, x, y, Emu(cw), Inches(4.05), fill=WHITE, line=RULE)
    box(s, x, y, Emu(cw), Pt(2.6), fill=NAVY if i == 0 else RULE)
    text(s, x + Inches(0.18), y + Inches(0.2), Emu(cw) - Inches(0.36), Inches(0.2),
         chip.upper(), size=8.5, color=MUTED, font=MONO_FONT)
    text(s, x + Inches(0.18), y + Inches(0.46), Emu(cw) - Inches(0.36), Inches(0.55),
         name, size=14, color=INK, font=HEAD_FONT, bold=True)
    text(s, x + Inches(0.18), y + Inches(1.16), Emu(cw) - Inches(0.36), Inches(0.45),
         price, size=26, color=NAVY, font=HEAD_FONT, bold=True)
    text(s, x + Inches(0.18), y + Inches(1.63), Emu(cw) - Inches(0.36), Inches(0.24),
         per, size=9.5, color=MUTED)
    bullets(s, x + Inches(0.18), y + Inches(1.98), Emu(cw) - Inches(0.36), pros,
            size=10.5, h=Inches(1.1), dash_color=POS, space=Pt(5))
    if cons:
        bullets(s, x + Inches(0.18), y + Inches(1.98) + Inches(0.27) * len(pros) + Inches(0.16),
                Emu(cw) - Inches(0.36), cons, size=10.5, h=Inches(1.2),
                dash_color=NEG, space=Pt(5))

# ───────────────────────────── 4 · the deciding rate ─────────────────────────────
s, y = slide_frame("The rate that decides the outcome",
                   kicker="SECTION 2.1 · WHY NO MIGRATION TARGET REACHES ₹21,000",
                   lede="The estate is memory-heavy and CPU-light — 44 GiB resident against 1.35 cores in use. The deciding figure is the price of one permanently resident unit of capacity: the nearest equivalent of an OCI baseline node, 4 OCPU and 24 GB, held warm for a month.")
table(s, MARGIN, y + Inches(0.34), Inches(8.1),
      ["Provider", "Nearest equivalent shape", ">₹ / month", ">× OCI"],
      [["OCI Ampere A1 — current", "4 OCPU / 24 GB", ("4,421", {"num": True}), ("1.0×", {"num": True, "color": POS})],
       ["AWS Graviton, ap-south-1", "m7g.xlarge — 4 vCPU / 16 GiB", ("≈ 11,000", {"num": True}), ("2.5×", {"num": True, "color": NEG})],
       ["DigitalOcean, Bangalore", "General purpose — 4 vCPU / 16 GB", ("≈ 11,100", {"num": True}), ("2.5×", {"num": True, "color": NEG})],
       ["Azure VM, Central India", "D4ps_v5 — 4 vCPU / 16 GiB", ("≈ 12,200", {"num": True}), ("2.8×", {"num": True, "color": NEG})],
       ["Azure Container Apps — idle", "4 vCPU / 24 GiB warm", ("≈ 19,300", {"num": True}), ("4.4×", {"num": True, "color": NEG})],
       ["Google Cloud Run — always on", "4 vCPU / 24 GiB warm", ("≈ 41,300", {"num": True}), ("9.3×", {"num": True, "color": ALERT})]],
      col_w=[30, 34, 20, 16], row_h=Inches(0.36))
bx = MARGIN + Inches(8.4)
box(s, bx, y + Inches(0.34), CONTENT_W - Inches(8.4), Inches(2.66), fill=TINT, line=RULE)
box(s, bx, y + Inches(0.34), Pt(3.2), Inches(2.66), fill=NAVY)
text(s, bx + Inches(0.22), y + Inches(0.52), CONTENT_W - Inches(8.85), Inches(2.3),
     [[("Memory is the constraint.", {"bold": True, "color": INK, "size": 13})],
      [("OCI Ampere A1 prices memory at ", {}), ("₹86.82 per GB-month", {"bold": True, "color": INK}),
       (". Every alternative charges ₹300 to ₹700 for the same gigabyte.", {})],
      [("Compute is where rivals come closest, at 2–3×. Memory is where they are 5–10× dearer — and memory is what this estate consumes.", {})]],
     size=11.5)
box(s, MARGIN, y + Inches(3.02), CONTENT_W, Inches(0.78), fill=NAVY)
text(s, MARGIN + Inches(0.3), y + Inches(3.18), CONTENT_W - Inches(0.6), Inches(0.5),
     [[("Screening rule: ", {"bold": True, "color": WHITE}),
       ("if a provider cannot price memory below roughly ₹150 per GB-month, it cannot reach Plan A — and the rest of the modelling is wasted effort.",
        {"color": WHITE})]], size=13)

# ───────────────────────────── 5 · evidence ─────────────────────────────
s, y = slide_frame("What the estate actually does", kicker="SECTION 4 · EVIDENCE BASE",
                   lede="Measured from the live estate and the production database on 17 August 2026, with the cluster re-checked on 27 August.")
table(s, MARGIN, y + Inches(0.3), CONTENT_W,
      ["Measurement", ">Value", "Implication"],
      [["Assessment submissions, last 30 days", ("648 (~21/day)", {"num": True}), "Peak hour was 32 submissions — one every 112 seconds"],
       ["Hours with any activity", ("251 of 720 (34.9%)", {"num": True}), "Two thirds of the month is completely idle"],
       ["Kubernetes CPU in use", ("1.35 of 20 cores", {"num": True}), "Nodes run at 2–11% utilisation"],
       ["Memory in use", ("44 of 120 GB", {"num": True}), "Memory is the binding cost driver, not CPU"],
       ["Sum of pod CPU requests", ("~6.5 cores", {"num": True}), "fast-api alone reserves 5× the cluster's real usage"],
       ["Production database size", ("10 GB", {"num": True}), "Held by a managed service costing ₹21,870/month"],
       ["Open database connections", ("213", {"num": True}), "PgBouncer takes this to about 20"],
       ["Block and boot storage", ("~2,900 GB", {"num": True}), "One static-site VM carries a 500 GB boot volume"],
       ["Container images in registry", ("407 GB + 151 GB", {"num": True}), "No retention policy has ever run"],
       ["verify-frame cold vs warm", ("3.3 s → 8 ms", {"num": True}), "The exam path can never scale to zero — in any plan"]],
      col_w=[30, 20, 50], row_h=Inches(0.335))

# ───────────────────────────── 6 · phase 0 ─────────────────────────────
s, y = slide_frame("Phase 0 — already banked", kicker="SECTION 3 · COMPLETE, 17 AUGUST 2026")
box(s, MARGIN, y + Inches(0.1), Inches(3.5), Inches(1.5), fill=NAVY)
text(s, MARGIN + Inches(0.3), y + Inches(0.34), Inches(3.0), Inches(0.4),
     "₹44,000", size=38, color=WHITE, font=HEAD_FONT, bold=True)
text(s, MARGIN + Inches(0.3), y + Inches(1.02), Inches(3.0), Inches(0.3),
     "per month, removed", size=12.5, color=RGBColor(0xC7, 0xD9, 0x66))
text(s, MARGIN, y + Inches(1.86), Inches(3.5), Inches(1.6),
     "The database was cut over to pl-prod-pg16 on 31 July, but the previous PostgreSQL 14 system — two instances of 2 OCPU and 32 GB — kept running in parallel with nothing connected to it.",
     size=12)
cx = MARGIN + Inches(3.9)
text(s, cx, y + Inches(0.1), Inches(4.2), Inches(0.3), "VERIFIED BEFORE DELETION",
     size=9.5, color=NAVY2, font=MONO_FONT)
bullets(s, cx, y + Inches(0.45), Inches(4.2), [
    "Zero client connections; 2 inserts, 3 updates, 0 deletes over eight days",
    "Last row written 3 August versus 17 August on the live system",
    "Diverged by 295 students and 5,629 assignment records — a frozen snapshot",
    "Zero references across all ten namespaces",
    "After deletion: all pods running, no new restarts",
], size=11.5)
cx2 = MARGIN + Inches(8.5)
text(s, cx2, y + Inches(0.1), Inches(4.2), Inches(0.3), "ARCHIVED, BOTH SURVIVING",
     size=9.5, color=NAVY2, font=MONO_FONT)
bullets(s, cx2, y + Inches(0.45), Inches(4.2), [
    "Logical dump — 2.17 GB, verified at 2,603 objects",
    "Manual system backup — 1 GB, 35-day retention",
    "Side-fix: three read-only helper scripts had been returning stale production data since 3 August; all repointed",
], size=11.5)

# ───────────────────────────── 7 · plan A shape ─────────────────────────────
s, y = slide_frame("Plan A — what changes", kicker="SECTION 5 · RE-SHAPE ON ORACLE CLOUD",
                   lede="Everything stays on OCI, in Mumbai. The saving comes from matching capacity to measured load, and switching capacity off when it is genuinely not needed.")
text(s, MARGIN, y + Inches(0.3), Inches(6.2), Inches(0.3), "TARGET ARCHITECTURE",
     size=9.5, color=NAVY2, font=MONO_FONT)
bullets(s, MARGIN, y + Inches(0.66), Inches(6.2), [
    [("Production — ", {"bold": True, "color": INK}), ("OKE Basic (free control plane), two A1 4/24 baseline nodes, one preemptible calendar burst node, self-managed PostgreSQL 16 with PgBouncer", {})],
    [("Frontends — ", {"bold": True, "color": INK}), ("Object Storage behind the existing load balancer", {})],
    [("DEV and UAT — ", {"bold": True, "color": INK}), ("one VM each, database merged in, scheduled off nights and weekends", {})],
    [("Batch — ", {"bold": True, "color": INK}), ("parsers and normalisation on Container Instances, billed per second", {})],
], size=12.5)
cx = MARGIN + Inches(6.7)
text(s, cx, y + Inches(0.3), Inches(6.0), Inches(0.3), "THREE RULES BAKED INTO THE DESIGN",
     size=9.5, color=NAVY2, font=MONO_FONT)
rules = [("The exam path never sleeps", "auth-node, student-node and fast-api stay warm during exam hours — 3.3 s cold against 8 ms warm is why."),
         ("App and database stay together", "Documented N+1 patterns turn a 2 ms hop into 100 ms+ per page."),
         ("Capacity leads demand", "Scale from the assessment schedule 20 minutes ahead, not from CPU metrics after the fact. Zero cold start.")]
ry = y + Inches(0.66)
for i, (t, d) in enumerate(rules):
    box(s, cx, ry, Inches(6.0), Inches(1.0), fill=TINT, line=RULE)
    text(s, cx + Inches(0.18), ry + Inches(0.13), Inches(0.5), Inches(0.24),
         f"RULE {i+1}", size=8.5, color=NAVY2, font=MONO_FONT)
    text(s, cx + Inches(0.9), ry + Inches(0.11), Inches(5.0), Inches(0.26), t,
         size=12.5, color=INK, bold=True)
    text(s, cx + Inches(0.18), ry + Inches(0.44), Inches(5.65), Inches(0.5), d, size=11)
    ry = ry + Inches(1.12)

# ───────────────────────────── 8 · plan A schedule ─────────────────────────────
s, y = slide_frame("Plan A — seven weeks, ₹84,800 → ₹21,000", kicker="SECTION 5 · DELIVERY SCHEDULE")
table(s, MARGIN, y + Inches(0.22), CONTENT_W,
      ["When", "Work", ">Saving", ">Run-rate after"],
      [["Week 0", "Raise pg16 backup retention 1 → 30 days. Resource Scheduler on DEV, UAT and the builder: off 20:00, on 08:00, weekdays", ("₹8,488", {"num": True, "color": POS}), ("₹76,300", {"num": True})],
       ["Week 1", "Merge DEV and UAT database VMs into their app hosts · registry retention · block volumes VPU 10 → 0 · 7-day baseline", ("₹10,300", {"num": True, "color": POS}), ("₹66,000", {"num": True})],
       ["Weeks 2–3", "Right-size pod requests to p95 + 20% · PodDisruptionBudgets · KEDA in report-only mode", ("enabler", {"num": True, "color": MUTED}), ("₹66,000", {"num": True})],
       ["Week 4", "New OKE Basic cluster, 2-node pool on 50 GB boots · drain, DNS switch · static sites to Object Storage", ("₹21,500", {"num": True, "color": POS}), ("₹44,500", {"num": True})],
       ["Weeks 5–6", "PgBouncer (213 → ~20 connections) · self-managed PG16 · timed restore test · cutover · 7-day soak", ("₹17,449", {"num": True, "color": POS}), ("₹27,000", {"num": True})],
       ["Week 7", "Calendar scaler — burst node warm 20 min before each window · KEDA scale-to-zero on idle workers", ("₹3–6,000", {"num": True, "color": POS}), ("≈ ₹21,000", {"num": True})]],
      col_w=[11, 55, 16, 18], row_h=Inches(0.58), font=11)
text(s, MARGIN, y + Inches(4.28), CONTENT_W, Inches(0.3),
     [[("Weeks 0 and 1 are unconditional — ", {"bold": True, "color": INK}),
       ("they apply whichever option is chosen, and none of them touches the production data path.", {})]], size=12.5)

# ───────────────────────────── 9 · plan A build-up + environments ─────────────
s, y = slide_frame("Plan A — where the ₹21,000 goes", kicker="SECTION 5.5 · COST BUILD-UP AND ENVIRONMENTS")
table(s, MARGIN, y + Inches(0.22), Inches(7.0),
      ["Component", "Shape", ">₹ / month"],
      [["Baseline OKE nodes (2)", "A1 4 OCPU / 24 GB each", ("8,842", {"num": True})],
       ["Calendar burst node", "preemptible, ~40 h/month", ("200", {"num": True})],
       ["Self-managed PostgreSQL", "A1 4/24 plus 100 GB", ("4,421", {"num": True})],
       ["DEV + UAT + builder", "3 VMs, scheduled 36% of hours", ("4,774", {"num": True})],
       ["OKE Basic control plane", "—", ("0", {"num": True})],
       ["Frontends", "Object Storage + existing LB", ("0", {"num": True})],
       ["Block storage + performance units", "~950 GB, VPU 0 except DB", ("2,620", {"num": True})],
       ["Images, storage, email, DNS, checks", "—", ("1,083", {"num": True})],
       [("Plan A total", {"bold": True}), ("±5%", {}), ("≈ 21,940", {"num": True, "bold": True})]],
      col_w=[40, 40, 20], row_h=Inches(0.335))
cx = MARGIN + Inches(7.35)
text(s, cx, y + Inches(0.22), Inches(4.8), Inches(0.3), "BY ENVIRONMENT",
     size=9.5, color=NAVY2, font=MONO_FONT)
envs = [("DEV", "₹6,631 → ≈₹2,200", "Merge the database VM into the app host; schedule off nights and weekends with a one-click override"),
        ("UAT", "₹7,673 → ≈₹2,500", "Same consolidation and schedule, with an override for planned test windows"),
        ("PROD", "₹70,500 → ≈₹16,300", "Right-size requests → OKE Basic with 2 nodes + burst → static sites to Object Storage → self-managed PG16")]
ey = y + Inches(0.6)
for name, delta, desc in envs:
    box(s, cx, ey, Inches(4.9), Inches(1.15), fill=WHITE, line=RULE)
    box(s, cx, ey, Pt(2.6), Inches(1.15), fill=NAVY)
    text(s, cx + Inches(0.2), ey + Inches(0.12), Inches(1.2), Inches(0.26), name,
         size=13, color=NAVY, font=HEAD_FONT, bold=True)
    text(s, cx + Inches(1.35), ey + Inches(0.14), Inches(3.4), Inches(0.26), delta,
         size=12, color=POS, bold=True, font=MONO_FONT)
    text(s, cx + Inches(0.2), ey + Inches(0.45), Inches(4.5), Inches(0.6), desc, size=10.5)
    ey = ey + Inches(1.28)

# ───────────────────────────── 10 · prerequisites ─────────────────────────────
s, y = slide_frame("Four gaps gate the schedule", kicker="SECTION 5.3 · INSTRUMENTATION PREREQUISITES",
                   lede="Verified on the live cluster, 27 August. Prometheus, kube-state-metrics, node-exporter, Grafana, Loki, Promtail and metrics-server are already deployed in production — no new monitoring platform is needed. Four gaps remain, together costing ₹100–200 per month.")
table(s, MARGIN, y + Inches(0.26), CONTENT_W,
      ["Prerequisite", "Current state", ">Gates", "Why it is required", ">Cost"],
      [["Alertmanager in production", ("Absent — UAT has it", {"color": ALERT, "bold": False}), ("Week 2", {"num": True}), "Production drops from five worker nodes to two plus a burst node. Less headroom demands alerting, not less of it", ("≈ ₹0", {"num": True})],
       ["Prometheus retention 7 → 30 days", ("7 days", {"color": ALERT, "bold": False}), ("Week 2", {"num": True}), "Week 1 captures a 7-day baseline and weeks 2–3 right-size to p95 + 20%; at 7 days the evidence expires first", ("≈ ₹70", {"num": True})],
       ["PodDisruptionBudgets", ("None defined", {"color": ALERT, "bold": False}), ("Week 4", {"num": True}), "Week 4 drains nodes one at a time; without a budget a drain can remove the last replica of a customer-facing API", ("₹0", {"num": True})],
       ["postgres_exporter", ("Absent", {"color": ALERT, "bold": False}), ("Weeks 5–6", {"num": True}), "D-1 transfers backups and recovery to the team; connections, replication lag and WAL success must be visible first", ("≈ ₹0", {"num": True})]],
      col_w=[22, 16, 10, 42, 10], row_h=Inches(0.62), font=10.5)
text(s, MARGIN, y + Inches(3.22), CONTENT_W, Inches(0.5),
     [[("Not required: ", {"bold": True, "color": INK}),
       ("Grafana Alloy replaces a working Promtail and changes no cost line. OpenLens and Headlamp duplicate Grafana and kubectl. Terraform is worth capturing at the week-4 rebuild, not as a mid-migration import.", {})]], size=12)

# ───────────────────────────── 11 · plan B ─────────────────────────────
s, y = slide_frame("Plan B — Google Cloud Run, Mumbai", kicker="SECTION 7 · ALTERNATIVE",
                   lede="Request-driven containers, a managed database, no cluster to operate. The price depends entirely on whether the queue engine is also rewritten.")
table(s, MARGIN, y + Inches(0.3), Inches(6.4),
      ["B1 — lift and shift", ">₹ / month"],
      [["Always-warm tier — auth, student, fast-api", ("24,600", {"num": True})],
       ["Queue workers held always-on", ("11,500", {"num": True})],
       ["Scale-to-zero tier", ("700", {"num": True})],
       ["Cloud SQL PostgreSQL — 2 vCPU / 8 GB", ("10,000", {"num": True})],
       ["Memorystore Redis — 1 GB", ("2,000", {"num": True})],
       ["Storage, registry, LB, egress, logging", ("4,000", {"num": True})],
       ["DEV + UAT — true scale-to-zero", ("2,000", {"num": True})],
       [("B1 total", {"bold": True}), ("≈ 44,000", {"num": True, "bold": True})]],
      col_w=[74, 26], row_h=Inches(0.335))
cx = MARGIN + Inches(6.75)
table(s, cx, y + Inches(0.3), Inches(5.35),
      ["B2 — after the queue rewrite", ">₹ / month"],
      [["Workers become request-driven via Pub/Sub push", ("−11,500", {"num": True, "color": POS})],
       ["fast-api warm only in scheduled windows", ("−6,600", {"num": True, "color": POS})],
       ["All other lines unchanged", ("—", {"num": True})],
       [("B2 total", {"bold": True}), ("≈ 26,000", {"num": True, "bold": True})]],
      col_w=[74, 26], row_h=Inches(0.42))
box(s, cx, y + Inches(2.32), Inches(5.35), Inches(1.62), fill=TINT, line=RULE)
box(s, cx, y + Inches(2.32), Pt(3.2), Inches(1.62), fill=ALERT)
text(s, cx + Inches(0.22), y + Inches(2.5), Inches(4.9), Inches(1.3),
     [[("The blocker that sets the price: BullMQ", {"bold": True, "color": INK, "size": 12.5})],
      [("Cloud Run bills a container only while it handles a request. A worker blocking on Redis is not handling one, so it needs CPU always allocated — a virtual machine at roughly twelve times the price. Escaping it means re-deriving retry, back-off and delayed-job semantics.", {})]],
     size=10.5)
text(s, MARGIN, y + Inches(3.46), Inches(6.4), Inches(0.9),
     [[("Migration: 10 to 14 weeks. ", {"bold": True, "color": INK}),
       ("The high-risk stages are the data platform — 10 GB by logical replication, 160 GB of media re-signed — and the queue rewrite. Rollback after the data moves is difficult.", {})]], size=11.5)

# ───────────────────────────── 12 · plan C ─────────────────────────────
s, y = slide_frame("Plan C — Azure Container Apps, Central India", kicker="SECTION 8 · ALTERNATIVE, SELF-MANAGED DATABASE",
                   lede="Costed against the live inventory: 20 API deployments, 12 frontends, four queue workers, measured resident memory per service.")
table(s, MARGIN, y + Inches(0.3), Inches(6.9),
      ["Component", ">₹ / month"],
      [["Container Apps — always-warm APIs", ("10,600", {"num": True})],
       ["Container Apps — queue workers (KEDA to zero)", ("800", {"num": True})],
       ["Container Apps — scale-to-zero tier, 12 services", ("5,500", {"num": True})],
       ["PostgreSQL host — D2pds_v5, self-managed", ("5,550", {"num": True})],
       ["Database storage + backup target", ("2,130", {"num": True})],
       ["Redis Basic C1 · frontends · registry", ("8,550", {"num": True})],
       ["Log Analytics · media, egress, DNS, Key Vault", ("3,900", {"num": True})],
       ["DEV and UAT", ("4,300", {"num": True})],
       [("Plan C total", {"bold": True}), ("≈ 41,300", {"num": True, "bold": True})]],
      col_w=[76, 24], row_h=Inches(0.335))
cx = MARGIN + Inches(7.25)
text(s, cx, y + Inches(0.3), Inches(4.9), Inches(0.3), "WHERE IT BEATS CLOUD RUN",
     size=9.5, color=NAVY2, font=MONO_FONT)
bullets(s, cx, y + Inches(0.66), Inches(4.9), [
    [("Idle billing. ", {"bold": True, "color": INK}), ("A warm replica bills at roughly one-eighth the active rate; Cloud Run charges full price.", {})],
    [("KEDA is built in. ", {"bold": True, "color": INK}), ("BullMQ workers scale to zero with no queue rewrite — B2's saving without B2's risk.", {})],
], size=11.5)
text(s, cx, y + Inches(1.98), Inches(4.9), Inches(0.3), "WHERE IT STILL LOSES",
     size=9.5, color=NEG, font=MONO_FONT)
bullets(s, cx, y + Inches(2.34), Inches(4.9), [
    [("Memory bills flat at ₹690 per GiB", {"bold": True, "color": INK}), (" whether idle or active; 44 GiB resident sets a floor near ₹30,000.", {})],
    [("Self-managing saves only ₹5,300 here", {"bold": True, "color": INK}), (" against ₹17,449 under Plan A — and removes the “nothing to operate” benefit.", {})],
    "Same migration exposure as Plan B; 9–12 weeks; difficult rollback.",
], size=11.5, dash_color=NEG)
text(s, MARGIN, y + Inches(3.78), Inches(6.9), Inches(0.4),
     [[("Lean variant ≈ ₹27,000 ", {"bold": True, "color": INK}),
       ("with reserved compute, Static Web Apps Free, Redis C0, Basic Logs and cron-warmed fast-api.", {})]], size=11.5)

# ───────────────────────────── 13 · providers screened ─────────────────────────
s, y = slide_frame("Every other provider screened", kicker="SECTION 9.1–9.2 · PLATFORM AND REGION FEASIBILITY")
table(s, MARGIN, y + Inches(0.22), CONTENT_W,
      ["Platform", "India region", "Assessment"],
      [["Alibaba Cloud", ("Withdrawn", {"color": ALERT, "bold": False}), "India data centre closed in 2022–23. Fails the India requirement outright"],
       ["Tencent Cloud", "Mumbai", "Rejected on procurement and sovereignty grounds, not price — student personal data on a Chinese-owned cloud will fail customer security review"],
       ["IBM Cloud", ("Chennai classic only", {"color": NEG, "bold": False}), "No full multi-zone India region for the managed services needed; no Ampere-class pricing; thin local ecosystem"],
       ["DigitalOcean", "Bangalore only", "Compute competitive at ~2.5× OCI, but memory near ₹700/GB-month — the memory bill alone would exceed ₹30,000. No in-country failover"],
       ["Salesforce", "n/a", "Not an infrastructure provider — SaaS plus Heroku, priced per dyno above Cloud Run"],
       ["Render · Neon", ("None", {"color": ALERT, "bold": False}), "No India region at all; two options from the original proof-of-concept stack fail requirement 1"],
       ["AWS · Azure · GCP", "Mumbai / Central India", "All viable on region. Costed as Plans B and C, or rejected on the rate in section 2.1 — between them and OCI sits a 2.5× to 9.3× gap"]],
      col_w=[16, 18, 66], row_h=Inches(0.56), font=11)

# ───────────────────────────── 14 · vector databases ─────────────────────────
s, y = slide_frame("Vector databases cannot replace PostgreSQL", kicker="SECTION 9.3 · DATA-STORE FEASIBILITY",
                   lede="Measured in production on 27 August: about 104,000 embeddings in 1.7 GB, inside an 11 GB database that is otherwise students, assessments, assignments and scores.")
table(s, MARGIN, y + Inches(0.42), Inches(5.6),
      ["Table", ">Size", ">Rows"],
      [["institute.institutes_campuses", ("1,380 MB", {"num": True}), ("86,903", {"num": True})],
       ["admin.mongo_db_cities", ("91 MB", {"num": True}), ("5,736", {"num": True})],
       ["institute.degrees", ("70 MB", {"num": True}), ("4,454", {"num": True})],
       ["corporate.job_roles", ("69 MB", {"num": True}), ("2,359", {"num": True})],
       ["institute.streams", ("64 MB", {"num": True}), ("4,078", {"num": True})],
       ["Three smaller tables", ("≈ 5 MB", {"num": True}), ("≈ 300", {"num": True})],
       [("Total", {"bold": True}), ("≈ 1.7 GB", {"num": True, "bold": True}), ("≈ 104,000", {"num": True, "bold": True})]],
      col_w=[54, 23, 23], row_h=Inches(0.335), font=10.5)
cx = MARGIN + Inches(5.95)
bullets(s, cx, y + Inches(0.42), Inches(6.2), [
    [("No joins, no foreign keys, no transactions. ", {"bold": True, "color": INK}), ("They are similarity engines over embeddings — an additional store, never a substitute for the system of record.", {})],
    [("Cost. ", {"bold": True, "color": INK}), ("pgvector on the existing database costs nothing. Pinecone, Qdrant and Weaviate each start near ₹2,200–4,400 per month — 10–20% of the entire Plan A target — to serve 104,000 vectors.", {})],
    [("Region. ", {"bold": True, "color": INK}), ("Pinecone serverless has no India region, which fails requirement 1 outright.", {})],
    [("Correctness. ", {"bold": True, "color": INK}), ("Ranking is hybrid — taxonomy gate, trigram and vector fused in one query. Splitting the vectors out removes the family gate that fixed “nurse resembles accountant”.", {})],
    [("Direction of travel. ", {"bold": True, "color": INK}), ("Elasticsearch was just consolidated into PostgreSQL. Adding a vector store puts a data store back.", {})],
], size=11.5)
text(s, MARGIN, y + Inches(3.22), Inches(5.6), Inches(1.0),
     [[("Revisit at roughly ten million vectors — the estate is two orders of magnitude away.", {"bold": True, "color": INK})],
      [("A cheaper win exists now: institutes_campuses holds ~16 KB per row, so halfvec or binary quantisation would shrink it with no new service.", {})]], size=11)

# ───────────────────────────── 15 · decisions ─────────────────────────────
s, y = slide_frame("Decisions to sign off", kicker="SECTION 10 · DECISION REGISTER")
table(s, MARGIN, y + Inches(0.22), CONTENT_W,
      ["Ref", "Decision", "Owner", ">If accepted", ">If declined"],
      [["D-0", "Adopt Plan A, Plan B or Plan C", "CTO", ("A → ₹21,000", {"num": True, "color": POS}), ("B → ₹26–44k · C → ₹27–41.3k", {"num": True})],
       ["D-1", "Self-manage PostgreSQL — team owns backups and recovery for a 10 GB database", "CTO", ("−₹17,449", {"num": True, "color": POS}), ("floor becomes ₹38,000", {"num": True})],
       ["D-2", "Two baseline nodes with a calendar burst node, not three always on", "Infra", ("−₹4,421", {"num": True, "color": POS}), ("₹25,400", {"num": True})],
       ["D-3", "Migrate to an OKE Basic cluster — new cluster, DNS cutover, one-way", "Infra", ("−₹6,834", {"num": True, "color": POS}), ("₹27,800", {"num": True})],
       ["D-4", "DEV, UAT and builder off nights and weekends, with a wake-up override", "Eng + QA", ("−₹8,488", {"num": True, "color": POS}), ("₹29,500", {"num": True})],
       ["D-5", "Rewrite BullMQ to Pub/Sub push — Plan B only; unnecessary under Plan C", "Backend", ("−₹11,500", {"num": True, "color": POS}), ("B remains ₹44,000", {"num": True})],
       ["D-6", "Move embeddings to a dedicated vector database — not recommended", "CTO", ("+₹2,200–4,400", {"num": True, "color": NEG}), ("₹0 — pgvector retained", {"num": True, "color": POS})]],
      col_w=[6, 40, 11, 19, 24], row_h=Inches(0.5), font=11)

# ───────────────────────────── 16 · risks and guardrails ─────────────────────
s, y = slide_frame("Risks and guardrails", kicker="SECTIONS 11–12 · RISK REGISTER AND ROLLBACK")
text(s, MARGIN, y + Inches(0.2), Inches(6.3), Inches(0.3), "OPEN RISKS THAT GATE THE PLAN",
     size=9.5, color=NEG, font=MONO_FONT)
risks = [("HIGH", "pg16 backup retention is 1 day", "The only production database has a 24-hour restore window. Raise to 30 days before anything else"),
         ("HIGH", "Credentials in plain ConfigMaps", "Cloud, AI and database credentials readable by any cluster user. Move to Secrets or OCI Vault"),
         ("HIGH", "Overnight CI depends on scheduled hosts", "D-4 stops DEV, UAT and the builder at 20:00. Decide wake-on-demand or move jobs into working hours"),
         ("MED", "Backup success is not alerted", "Alert on backup age and WAL archiving failure; a runbook is a D-1 deliverable"),
         ("MED", "Self-managed database ownership", "Mitigated by a completed restore drill before cutover, not after")]
ry = y + Inches(0.56)
for sev, t, d in risks:
    c = ALERT if sev == "HIGH" else NEG
    box(s, MARGIN, ry + Inches(0.04), Inches(0.52), Inches(0.22), line=c)
    text(s, MARGIN, ry + Inches(0.06), Inches(0.52), Inches(0.2), sev, size=8,
         color=c, font=MONO_FONT, align=PP_ALIGN.CENTER)
    text(s, MARGIN + Inches(0.66), ry, Inches(5.6), Inches(0.24), t, size=12,
         color=INK, bold=True)
    text(s, MARGIN + Inches(0.66), ry + Inches(0.26), Inches(5.6), Inches(0.4), d, size=10.5)
    ry = ry + Inches(0.82)
cx = MARGIN + Inches(6.9)
box(s, cx, y + Inches(0.2), Inches(5.25), Inches(4.28), fill=NAVY)
text(s, cx + Inches(0.28), y + Inches(0.42), Inches(4.7), Inches(0.3),
     "ACCEPTANCE CRITERIA FOR EVERY PRODUCTION CHANGE", size=9.5,
     color=RGBColor(0xC7, 0xD9, 0x66), font=MONO_FONT)
bullets(s, cx + Inches(0.28), y + Inches(0.82), Inches(4.7),
        ["API p95 latency no more than 10% worse than the week-1 baseline",
         "Error and timeout rates do not increase",
         "No sustained node CPU above 70–75% after right-sizing",
         "No OOM kills, memory pressure or unexpected eviction",
         "Database connections below 70% of the configured limit",
         "Queue delay and processing time do not materially increase",
         "Login, assessment, scoring, proctoring and dashboard smoke tests pass",
         "Every capacity change has a written, tested rollback before it starts"],
        size=11, h=Inches(2.5), dash_color=RGBColor(0xC7, 0xD9, 0x66), space=Pt(4))
for _p in s.shapes[-1].text_frame.paragraphs:
    for _r in _p.runs[1:]:
        _r.font.color.rgb = RGBColor(0xED, 0xF0, 0xE4)
text(s, cx + Inches(0.28), y + Inches(3.68), Inches(4.7), Inches(0.7),
     "Observation windows: 48 hours after a node drain · 7 days after the database cutover · 48 hours before deleting any replaced resource.",
     size=10.5, color=RGBColor(0xE3, 0xE7, 0xDC))

# ───────────────────────────── 17 · recommendation ─────────────────────────
s, y = slide_frame("Recommendation — approve Plan A", kicker="SECTION 14 · STATED LAST, AFTER THE EVIDENCE")
box(s, MARGIN, y + Inches(0.10), CONTENT_W, Inches(1.05), fill=TINT, line=RULE)
box(s, MARGIN, y + Inches(0.10), Pt(3.2), Inches(1.05), fill=NAVY)
text(s, MARGIN + Inches(0.28), y + Inches(0.26), CONTENT_W - Inches(0.6), Inches(0.85),
     [[("Plan A costs about half of the best achievable Cloud Run outcome, delivers ₹8,488 in week 0 rather than week 10, needs no data migration, and is reversible at every step. Plan B2 reaches near-parity only after a queue rewrite; Plan C is the strongest migration option on risk, yet memory at ₹690 per GiB against 44 GiB resident sets a floor it cannot clear.", {})]],
     size=13)
table(s, MARGIN, y + Inches(1.32), CONTENT_W,
      ["Sequence", "Action", "Owner", ">Effect"],
      [["Immediately", "Raise pl-prod-pg16 backup retention 1 → 30 days. Unconditional; it gates the database work under every option", "Infra", ("risk fix", {"num": True, "color": MUTED})],
       ["Week 0", "Resource Scheduler on DEV, UAT and the builder, with a wake-up override, after resolving the overnight CI dependency", "Infra · Eng · QA", ("₹8,488", {"num": True, "color": POS})],
       ["Week 1", "Merge the DEV and UAT database hosts, registry retention, block volumes, seven-day baseline", "Infra", ("₹10,300", {"num": True, "color": POS})],
       ["Before week 2", "Close the instrumentation gaps — Alertmanager in production and 30-day metric retention", "Infra", ("enabler", {"num": True, "color": MUTED})],
       ["Weeks 2–7", "Right-sizing, the OKE Basic rebuild, static sites, then the database cutover — each behind the guardrails", "Infra · Backend · DBA", ("₹45,000", {"num": True, "color": POS})]],
      col_w=[12, 52, 18, 18], row_h=Inches(0.46), font=11)
box(s, MARGIN, y + Inches(4.02), CONTENT_W, Inches(0.72), fill=NAVY)
text(s, MARGIN + Inches(0.28), y + Inches(4.14), CONTENT_W - Inches(0.56), Inches(0.5),
     [[("What is being asked: ", {"bold": True, "color": WHITE}),
       ("decision D-0 — adopt Plan A. Weeks 0 and 1 are unconditional and can proceed while D-1 to D-4 are considered.", {"color": WHITE})]], size=12.5)

out = "/home/ubuntu/pluginlive-kb/Infrastructure/pluginlive-infrastructure-cost-plans.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
